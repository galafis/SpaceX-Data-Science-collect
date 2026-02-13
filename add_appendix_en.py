from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from slide_content_en import appendix_charts

SCRIPT_DIR = Path(__file__).resolve().parent

# --- Paths ---
presentation_path = str(SCRIPT_DIR / "spacex_presentation_watermarked_en.pptx")
output_path = str(SCRIPT_DIR / "spacex_presentation_appendix_en.pptx")


def add_caption(slide, text, left, top, width, height):
    """Add a centered caption text box to a slide."""
    try:
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = text
        p = tf.paragraphs[0]
        p.font.size = Pt(10)
        p.alignment = PP_ALIGN.CENTER
    except Exception as e:
        print(f"Error adding caption '{text}': {e}")


# --- Main script ---
try:
    prs = Presentation(presentation_path)
    print(f"Loaded presentation: {presentation_path}")

    # --- Find Appendix Title Slide Index ---
    apendice_title_idx = -1
    for i, slide in enumerate(prs.slides):
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            if "appendix" in slide.shapes.title.text.lower():
                apendice_title_idx = i
                print(f"Found Appendix title slide at index {i}")
                break

    if apendice_title_idx == -1:
        print("Warning: Appendix title slide not found. Adding charts at the end.")

    # --- Add Appendix Chart Slides ---
    try:
        content_layout = prs.slide_layouts[1]
    except IndexError:
        print("Error: Title and Content layout (index 1) not found. Using layout 0.")
        content_layout = prs.slide_layouts[0]

    for chart_info in appendix_charts:
        try:
            slide = prs.slides.add_slide(content_layout)
            chart_title = chart_info.get("title", "Untitled Chart")
            print(f"Added new slide for appendix chart: {chart_title}")

            # Add Title
            title_shape = slide.shapes.title
            if title_shape:
                title_shape.text = chart_info.get("title", "Appendix Chart")
            else:
                print("Warning: Title placeholder not found on appendix chart slide.")

            # Add Image
            img_path = chart_info.get("image_path")
            if img_path:
                full_img_path = str(SCRIPT_DIR / img_path)
                left = Inches(1.0)
                top = Inches(1.5)
                height = Inches(4.5)
                try:
                    pic = slide.shapes.add_picture(full_img_path, left, top, height=height)
                    print(f"Added image: {full_img_path}")

                    caption_text = chart_info.get("caption", "")
                    if caption_text:
                        caption_top = top + height + Inches(0.1)
                        caption_width = pic.width
                        caption_height = Inches(0.5)
                        add_caption(slide, caption_text, left, caption_top, caption_width, caption_height)
                        print(f"Added caption: {caption_text}")

                except FileNotFoundError:
                    print(f"Error: Image file not found at {full_img_path}")
                except Exception as e_img:
                    print(f"Error adding image {full_img_path}: {e_img}")
            else:
                print("Warning: No image path specified for this appendix chart.")

        except Exception as e_slide:
            chart_title_error = chart_info.get("title", "Unknown Chart")
            print(f"Error adding slide for chart '{chart_title_error}': {e_slide}")

    # --- Save the final presentation ---
    prs.save(output_path)
    print(f"Presentation with appendix (English) saved to: {output_path}")

except Exception as e:
    print(f"An error occurred during appendix addition: {e}")
