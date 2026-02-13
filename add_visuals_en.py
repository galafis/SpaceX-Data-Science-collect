import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

SCRIPT_DIR = Path(__file__).resolve().parent

# --- Paths ---
presentation_path = str(SCRIPT_DIR / "spacex_presentation_draft_en.pptx")
output_path = str(SCRIPT_DIR / "spacex_presentation_visuals_en.pptx")

# Image paths
img_config_chart = str(SCRIPT_DIR / "rocket_configurations_chart_en.png")
img_launch_sites = str(SCRIPT_DIR / "launch_sites_chart_en.png")
img_outcomes_chart = str(SCRIPT_DIR / "launch_outcomes_chart_en.png")
img_landings_chart = str(SCRIPT_DIR / "booster_landings_chart_en.png")

# --- Helper function to find slide by title (approximate match) ---
def find_slide_by_title(prs, title_keyword):
    for i, slide in enumerate(prs.slides):
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            if title_keyword.lower() in slide.shapes.title.text.lower():
                print(f"Found slide '{slide.shapes.title.text}' (index {i}) matching keyword '{title_keyword}'")
                return slide
    print(f"Warning: Slide with title containing '{title_keyword}' not found.")
    return None

# --- Main script ---
try:
    prs = Presentation(presentation_path)
    print(f"Loaded presentation: {presentation_path}")

    # 1. Rocket Configurations Chart -> Falcon 9 Evolution slide
    evo_slide = find_slide_by_title(prs, "Evolution")
    if evo_slide:
        try:
            left = Inches(5.5)
            top = Inches(4.0)
            height = Inches(2.5)
            evo_slide.shapes.add_picture(img_config_chart, left, top, height=height)
            print(f"Added {img_config_chart} to Evolution slide.")
        except Exception as e:
            print(f"Error adding {img_config_chart}: {e}")

    # 2. Launch Outcomes & Sites Charts -> Launch Statistics slide
    stats_slide = find_slide_by_title(prs, "Statistics")
    if stats_slide:
        try:
            left1 = Inches(5.0)
            top1 = Inches(2.0)
            height1 = Inches(2.0)
            stats_slide.shapes.add_picture(img_outcomes_chart, left1, top1, height=height1)
            print(f"Added {img_outcomes_chart} to Statistics slide.")

            left2 = Inches(5.0)
            top2 = Inches(4.5)
            height2 = Inches(2.0)
            stats_slide.shapes.add_picture(img_launch_sites, left2, top2, height=height2)
            print(f"Added {img_launch_sites} to Statistics slide.")
        except Exception as e:
            print(f"Error adding charts to Statistics slide: {e}")

    # 3. Booster Landings Chart -> Reusability slide
    reuse_slide = find_slide_by_title(prs, "Reusability")
    if reuse_slide:
        try:
            left = Inches(5.5)
            top = Inches(4.0)
            height = Inches(2.5)
            reuse_slide.shapes.add_picture(img_landings_chart, left, top, height=height)
            print(f"Added {img_landings_chart} to Reusability slide.")
        except Exception as e:
            print(f"Error adding {img_landings_chart}: {e}")

    # --- Save the presentation with visuals ---
    prs.save(output_path)
    print(f"Presentation with visuals saved to: {output_path}")

except Exception as e:
    print(f"An error occurred during visual addition: {e}")
