from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = Path(__file__).resolve().parent

# Paths
input_pptx = str(SCRIPT_DIR / "spacex_presentation_appendix_en.pptx")
output_pptx = str(SCRIPT_DIR / "spacex_presentation_final_title_fixed_en.pptx")
author_name = "Gabriel Demetrios Lafis"
date_str = "May 2, 2025"

try:
    prs = Presentation(input_pptx)
    print(f"Loaded presentation: {input_pptx}")

    if len(prs.slides) > 0:
        slide1 = prs.slides[0]
        print("Accessing first slide to add author/date text box.")

        left = Inches(1.0)
        top = Inches(5.5)
        width = Inches(6.0)
        height = Inches(0.5)

        txBox = slide1.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"{author_name}\n{date_str}"

        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = "Calibri"
                run.font.size = Pt(14)

        print("Added text box with author and date to the first slide.")
    else:
        print("Warning: Presentation has no slides.")

    prs.save(output_pptx)
    print(f"Presentation with fixed title slide saved to: {output_pptx}")

except Exception as e:
    print(f"An error occurred while fixing the title slide: {e}")
