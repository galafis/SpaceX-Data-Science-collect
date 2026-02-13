from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

SCRIPT_DIR = Path(__file__).resolve().parent

# Paths
presentation_path = str(SCRIPT_DIR / "spacex_presentation_visuals_en.pptx")
output_path = str(SCRIPT_DIR / "spacex_presentation_watermarked_en.pptx")
watermark_text = "Gabriel Demetrios Lafis"

prs = Presentation(presentation_path)

# Define position and size for the watermark (bottom right corner)
width = prs.slide_width
height = prs.slide_height
wm_width = Inches(2.0)
wm_height = Inches(0.3)
wm_left = width - wm_width - Inches(0.2)
wm_top = height - wm_height - Inches(0.1)

# --- Add watermark to each slide ---
print("Adding watermark to each slide (English version)...")
slides_processed = 0
for i, slide in enumerate(prs.slides):
    try:
        txBox = slide.shapes.add_textbox(wm_left, wm_top, wm_width, wm_height)
        tf = txBox.text_frame

        tf.text = watermark_text
        tf.margin_bottom = Pt(0)
        tf.margin_top = Pt(0)
        tf.margin_left = Pt(0)
        tf.margin_right = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = False

        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT

        font = p.font
        font.size = Pt(8)
        font.color.rgb = RGBColor(192, 192, 192)
        font.name = 'Calibri'

        slides_processed += 1

    except Exception as e:
        print(f"Error adding watermark to slide {i+1}: {e}")

print(f"Watermark added to {slides_processed} slides.")

# Save the presentation with watermark
try:
    prs.save(output_path)
    print(f"Presentation with watermark saved to {output_path}")
except Exception as e:
    print(f"Error saving presentation with watermark: {e}")
