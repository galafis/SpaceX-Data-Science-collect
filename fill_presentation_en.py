import sys
sys.path.append("/opt/.manus/.sandbox-runtime")
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN # Needed for potential future formatting
import importlib.util

# --- Import slide content from the external file ---
try:
    spec = importlib.util.spec_from_file_location("slide_content_module", "/home/ubuntu/slide_content_en.py")
    slide_content_module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(slide_content_module)
    slide_content_en = slide_content_module.slide_content_en
    print("Successfully imported English slide content.")
except Exception as e:
    print(f"Error importing slide content: {e}")
    sys.exit(1) # Exit if content cannot be loaded

# --- Paths ---
template_path = '/home/ubuntu/upload/ds-capstone-template-coursera.pptx'
output_path = '/home/ubuntu/spacex_presentation_draft_en.pptx' # Intermediate file
author_name = "Gabriel Demetrios Lafis" # As requested by user

# --- Helper function to find placeholders ---
def find_placeholder(slide, placeholder_name):
    """Finds a placeholder shape by its name (e.g., 'TITLE', 'BODY', 'CONTENT')."""
    # Prioritize checking by explicit name if available in template
    for shape in slide.placeholders:
        # print(f"DEBUG: Checking placeholder: Name=	"{shape.name}	", Type={shape.placeholder_format.type}")
        if shape.name and placeholder_name.upper() in shape.name.upper():
             print(f"DEBUG: Found placeholder by name match: {shape.name}")
             return shape

    # Fallback to checking by type
    for shape in slide.placeholders:
        if (placeholder_name == 'TITLE' and shape.placeholder_format.type.name == 'TITLE') or \
           (placeholder_name == 'SUBTITLE' and shape.placeholder_format.type.name == 'SUBTITLE') or \
           (placeholder_name in ['BODY', 'CONTENT'] and shape.placeholder_format.type.name in ['BODY', 'CONTENT', 'OBJECT']):
            print(f"DEBUG: Found placeholder by type match: {shape.placeholder_format.type.name}")
            return shape

    # Fallback: Check slide.shapes.title (often works for title)
    if placeholder_name == 'TITLE' and slide.shapes.title:
        print("DEBUG: Found placeholder via slide.shapes.title")
        return slide.shapes.title

    # Fallback: Check by index (less reliable, use as last resort)
    if placeholder_name == 'SUBTITLE' and len(slide.placeholders) > 1:
         try:
             # Often index 1 is the subtitle or main content on title slides
             print(f"DEBUG: Attempting fallback to placeholder index 1 for {placeholder_name}")
             return slide.placeholders[1]
         except (KeyError, IndexError):
             pass # Ignore if index 1 doesn't exist
    if placeholder_name in ['BODY', 'CONTENT'] and len(slide.placeholders) > 1:
         try:
             # Often index 1 is the main content placeholder on content slides
             print(f"DEBUG: Attempting fallback to placeholder index 1 for {placeholder_name}")
             return slide.placeholders[1]
         except (KeyError, IndexError):
             pass # Ignore if index 1 doesn't exist

    print(f"Warning: Placeholder 	'{placeholder_name}	' not found reliably on slide.")
    return None

# --- Helper function to add slide with robust layout finding ---
def add_slide_robust(prs, layout_name_or_index, title_text, body_content):
    """Adds a slide using the specified layout, handling potential errors."""
    slide_layout = None
    try:
        if isinstance(layout_name_or_index, str):
            layout_names_to_try = [layout_name_or_index, layout_name_or_index.replace(" ", "_").upper()]
            for layout in prs.slide_layouts:
                if layout.name in layout_names_to_try:
                    slide_layout = layout
                    break
            if not slide_layout:
                 print(f"Warning: Layout 	'{layout_name_or_index}	' not found by name. Trying common indices.")
                 if layout_name_or_index == "Title Slide": slide_layout = prs.slide_layouts[0]
                 elif layout_name_or_index == "Title and Content": slide_layout = prs.slide_layouts[1]
                 elif layout_name_or_index == "Section Header": slide_layout = prs.slide_layouts[2]
                 else: slide_layout = prs.slide_layouts[1]
        elif isinstance(layout_name_or_index, int):
            slide_layout = prs.slide_layouts[layout_name_or_index]
        else:
             print(f"Warning: Invalid layout identifier 	'{layout_name_or_index}	'. Using default layout.")
             slide_layout = prs.slide_layouts[1]

        if not slide_layout:
             print(f"Error: Could not find or default to a suitable layout for 	'{title_text}	'. Skipping slide.")
             return None

        slide = prs.slides.add_slide(slide_layout)

        title_shape = find_placeholder(slide, 'TITLE')
        if title_shape:
            title_shape.text = title_text
        else:
            print(f"Warning: Title placeholder not found for slide 	'{title_text}	'.")

        if body_content:
            body_shape = find_placeholder(slide, 'CONTENT')
            if not body_shape:
                 body_shape = find_placeholder(slide, 'BODY')

            if body_shape:
                tf = body_shape.text_frame
                tf.clear()
                if isinstance(body_content, list):
                    for i, item in enumerate(body_content):
                        if i == 0:
                            tf.text = item
                        else:
                            p = tf.add_paragraph()
                            p.text = item
                            p.level = 0
                else:
                    tf.text = str(body_content)
            else:
                print(f"Warning: Body/Content placeholder not found for slide 	'{title_text}	'.")
        return slide

    except Exception as e:
        print(f"Error adding slide 	'{title_text}	': {e}")
        return None

# --- Main script ---
try:
    prs = Presentation(template_path)
    print(f"Loaded template: {template_path}")

    # --- Modify the FIRST slide (Title Slide) --- 
    if len(prs.slides) > 0:
        slide1 = prs.slides[0] # Get the first existing slide
        title_slide_content = None
        # Find the title slide content definition
        for content in slide_content_en:
            if content.get("layout_type") == "Title Slide":
                title_slide_content = content
                break
        
        if title_slide_content:
            print("Modifying existing first slide as Title Slide...")
            title = title_slide_content.get("title", "")
            subtitle = title_slide_content.get("subtitle", None)

            title_shape = find_placeholder(slide1, 'TITLE')
            subtitle_shape = find_placeholder(slide1, 'SUBTITLE') # Try finding SUBTITLE first

            if title_shape:
                title_shape.text = title
                print(f"Set title on first slide to: {title}")
            else:
                print("Warning: Title placeholder not found on first slide.")

            if subtitle_shape:
                 if subtitle:
                     subtitle_text = subtitle.replace("<Name>", author_name).replace("<Date>", "May 2, 2025")
                     subtitle_shape.text = subtitle_text
                     print(f"Set subtitle on first slide to: {subtitle_text}")
                 else:
                     default_subtitle = f"{author_name}\nMay 2, 2025"
                     subtitle_shape.text = default_subtitle
                     print(f"Set default subtitle on first slide to: {default_subtitle}")
            else:
                 print("Warning: Subtitle placeholder not found on first slide.")
        else:
            print("Warning: Title Slide content definition not found in slide_content_en.")

    else:
        print("Warning: Template presentation has no slides.")


    # --- Add OTHER slides based on the imported content --- 
    # Start adding from the second position if the first slide was modified
    start_index = 1 if len(prs.slides) > 0 else 0
    # We need to remove slides added previously by the old script if they exist beyond the original template slides
    # This is complex. A cleaner approach is to always start from the template and add slides.
    # Let's stick to adding slides after the potentially modified first slide.
    # However, the robust add_slide function adds slides at the end. 
    # For simplicity, let's clear existing slides beyond the first one if modifying the first.
    
    # --- REVISED APPROACH: Clear slides after the first, then add new ones ---
    if len(prs.slides) > 1:
        print(f"Clearing {len(prs.slides) - 1} slides after the first one before adding new content.")
        # Iterate backwards when removing
        for i in range(len(prs.slides) - 1, 0, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]

    # Now add all slides defined in slide_content_en EXCEPT the title slide
    for i, content in enumerate(slide_content_en):
        # Skip the title slide content as it was handled above
        if content.get("layout_type") == "Title Slide":
            continue

        print(f"Adding slide {i+1}: {content.get('title', 'Untitled')}")
        layout = content.get("layout_type", "Title and Content") # Default layout
        title = content.get("title", "")
        body = content.get("body", None)

        # Add other slides using the robust function (adds at the end)
        add_slide_robust(prs, layout, title, body)

    # --- Save the presentation ---
    prs.save(output_path) # output_path is '/home/ubuntu/spacex_presentation_draft_en.pptx'
    print(f"Presentation draft saved to: {output_path}")

except Exception as e:
    print(f"An error occurred during presentation creation: {e}")


